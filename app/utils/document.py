"""
Document processing utilities for text extraction, OCR, and file validation.
"""
import os
import re
import json
import numpy as np
import logging
import mimetypes
from typing import Optional
from PIL import Image, ImageEnhance, ImageOps

from .database import safe_db_query
from .pgvectorstore import get_vectorstore
from .time_provider import get_current_datetime
from .database import safe_db_query

# Import OCR and PDF processing dependencies
try:
    import pytesseract
    from PIL import Image, ImageEnhance, ImageOps, ImageFilter
except ImportError:
    logging.warning("OCR dependencies not available. Install with: pip install pytesseract pillow")
    pytesseract = None
    Image = None
    ImageEnhance = None
    ImageOps = None

try:
    import pdfplumber
    # Suppress noisy pdfminer warnings globally
    import warnings
    warnings.filterwarnings("ignore", category=UserWarning, module="pdfminer")
    warnings.filterwarnings("ignore", message=".*Cannot set.*color.*")
    warnings.filterwarnings("ignore", message=".*is an invalid float value.*")

    class _PdfMinerNoiseFilter(logging.Filter):
        def filter(self, record):
            message = record.getMessage()
            return "Cannot set gray non-stroke color" not in message
    
    for logger_name in ("pdfminer", "pdfminer.pdfinterp"):
        pdf_logger = logging.getLogger(logger_name)
        if not any(isinstance(f, _PdfMinerNoiseFilter) for f in pdf_logger.filters):
            pdf_logger.addFilter(_PdfMinerNoiseFilter())

except ImportError:
    logging.warning("PDF processing not available. Install with: pip install pdfplumber")
    pdfplumber = None

try:
    import fitz  # PyMuPDF
except ImportError:
    logging.warning("PyMuPDF not available. Install with: pip install PyMuPDF")
    fitz = None

try:
    import docx
except ImportError:
    logging.warning("python-docx not available. Install with: pip install python-docx")
    docx = None

try:
    import openpyxl
except ImportError:
    logging.warning("openpyxl not available. Install with: pip install openpyxl")
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    logging.warning("python-pptx not available. Install with: pip install python-pptx")
    Presentation = None

DEFAULT_TESSERACT_CONFIG = os.environ.get("TESSERACT_CONFIG", "--oem 3 --psm 3")
TESSERACT_CMD = os.environ.get("TESSERACT_CMD")

if pytesseract and TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
elif pytesseract and os.path.exists('/usr/bin/tesseract'):
    pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'

try:
    PDF_RENDER_SCALE = float(os.environ.get("PDF_RENDER_SCALE", "2.0"))
    if PDF_RENDER_SCALE <= 0:
        raise ValueError
except ValueError:
    PDF_RENDER_SCALE = 2.0

# def preprocess_image_for_ocr_old(image):
#     """Enhance an image to improve OCR accuracy using adaptive thresholding."""
#     if not Image:
#         return image

#     # 1. Konversi ke Grayscale (sudah benar)
#     img_gray = image.convert("L")

#     # 2. Resize jika terlalu kecil untuk detail yang lebih baik (sudah benar)
#     w, h = img_gray.size
#     if w < 1000 or h < 1000:
#         img_gray = img_gray.resize((w * 2, h * 2), Image.LANCZOS)

#     # 3. Konversi dari PIL Image ke format OpenCV
#     open_cv_image = np.array(img_gray)

#     # 4. Gunakan Thresholding Adaptif Gaussian
#     #    Ini adalah langkah kunci untuk mengatasi latar belakang yang bervariasi.
#     processed_img = cv2.adaptiveThreshold(
#         open_cv_image, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
#         cv2.THRESH_BINARY, 11, 2
#     )

#     # 5. Opsional: Sedikit denoise setelah thresholding
#     processed_img = cv2.medianBlur(processed_img, 3)

#     # 6. Konversi kembali ke PIL Image untuk Tesseract
#     return Image.fromarray(processed_img)

def preprocess_image_for_ocr(image):
    """Enhance an image to improve OCR accuracy."""
    if not Image:
        return image

    # Convert ke grayscale
    img = image.convert("L")

    # Autocontrast untuk menormalkan pencahayaan
    img = ImageOps.autocontrast(img)

    # Tingkatkan kontras & kecerahan sedikit
    img = ImageEnhance.Contrast(img).enhance(2.0)
    img = ImageEnhance.Brightness(img).enhance(1.1)

    # Kurangi noise dengan sedikit blur halus
    img = img.filter(ImageFilter.MedianFilter(size=3))

    # Sharpness untuk memperjelas tepi huruf
    img = ImageEnhance.Sharpness(img).enhance(2.0)

    # Resize jika gambar terlalu kecil
    w, h = img.size
    if w < 800:
        img = img.resize((w * 2, h * 2), Image.LANCZOS)

    # Binarize (thresholding) agar hasil OCR lebih stabil
    img = img.point(lambda x: 0 if x < 140 else 255, "1")

    return img

def clean_ocr_text(text):
    """Normalisasi hasil OCR agar lebih rapi."""
    # Hilangkan karakter non-alfanumerik yang aneh
    text = re.sub(r"[^a-zA-Z0-9.,()/%&+\-\s]", "", text)
    # Hilangkan spasi ganda
    text = re.sub(r"\s+", " ", text)
    # Format huruf besar pada awal kalimat
    text = text.strip()
    return text

def ocr_image(image):
    """Run Tesseract OCR on a PIL image with preprocessing."""
    if not (pytesseract and Image):
        return ''

    processed = preprocess_image_for_ocr(image)
    result = pytesseract.image_to_string(processed, lang="eng+ind", config=DEFAULT_TESSERACT_CONFIG).strip()
    return clean_ocr_text(result)

def perform_ocr_on_image(image, lang=None):
    """Run Tesseract OCR on a PIL image with preprocessing (legacy wrapper)."""
    return ocr_image(image)

def chunk_text(text, chunk_size=1536):
    """Split text into chunks of specified size."""
    chunks = []
    for i in range(0, len(text), chunk_size):
        chunks.append({"index": i // chunk_size, "content": text[i : i + chunk_size]})
    return chunks

def extract_text_from_pdf_image(path):
    """Extract text from PDF using PyMuPDF with OCR - matches grabber.py implementation."""
    if not (fitz and pytesseract and Image):
        return ''

    try:
        texts = []
        matrix = fitz.Matrix(PDF_RENDER_SCALE, PDF_RENDER_SCALE)
        with fitz.open(path) as doc:
            for page in doc:
                pix = page.get_pixmap(matrix=matrix)
                mode = "RGBA" if pix.alpha else "RGB"
                img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
                if pix.alpha:
                    img = img.convert("RGB")
                texts.append(ocr_image(img))
        return "\n".join(texts)
    except Exception as primary_error:
        # Fallback with temporary files like in grabber.py
        import shutil
        import uuid
        from pathlib import Path

        fallback_dir = Path("tmp") / f"{Path(path).stem}_ocr_{uuid.uuid4().hex}"
        fallback_texts = []
        try:
            fallback_dir.mkdir(parents=True, exist_ok=True)
            matrix = fitz.Matrix(PDF_RENDER_SCALE, PDF_RENDER_SCALE)
            with fitz.open(path) as doc:
                for index, page in enumerate(doc, start=1):
                    pix = page.get_pixmap(matrix=matrix)
                    img_path = fallback_dir / f"page_{index}.png"
                    pix.save(str(img_path))
                    with Image.open(img_path) as img:
                        fallback_texts.append(ocr_image(img))
        except Exception as fallback_error:
            raise fallback_error from primary_error
        finally:
            shutil.rmtree(fallback_dir, ignore_errors=True)

        return "\n".join(fallback_texts)

def data_path(*parts: str) -> str:
    """Return normalized path inside data directory.
    Handle various Docker deployment scenarios:
    - Production: /app (working dir) -> use /app/data
    - Dev container: /combiphar-be (working dir) -> use /combiphar-be/data
    - Legacy double path: /app/app -> use /app/data
    - Local: use ./data relative to current directory
    """
    cwd = os.getcwd()
    
    # Handle Docker production deployment
    if cwd == '/app' and os.path.exists('/app/data'):
        base = '/app/data'
    # Handle Docker dev deployment
    elif cwd == '/combiphar-be' and os.path.exists('/combiphar-be/data'):
        base = '/combiphar-be/data'
    # Handle legacy double path quirk
    elif cwd.endswith('/app/app') and os.path.exists('/app/data'):
        base = '/app/data'
    else:
        # Local development fallback
        base = os.path.abspath('./data')
    return os.path.join(base, *parts)

def verify_document_exists(stored_filename, original_filename=None):
    """
    Global helper function to verify if a document still exists in database and file system.
    Args:
        stored_filename (str): The stored filename (UUID-based) to verify
        original_filename (str, optional): The original filename for display
    Returns:
        dict: Dictionary containing existence status and metadata with keys:
            - exists_in_db (bool): Whether document exists in database
            - exists_in_file (bool): Whether document file exists in filesystem
            - original_filename (str): Original filename for display
            - stored_filename (str): Stored filename (UUID-based)
            - status (str): Overall status ('available', 'file_missing', 'db_missing', 'missing', 'error')
            - file_path (str): Full path to document file if found
            - storage_path (str): Storage path from database
    """
    result = {
        'exists_in_db': False,
        'exists_in_file': False,
        'original_filename': original_filename,
        'stored_filename': stored_filename,
        'status': 'missing',
        'file_path': None,
        'storage_path': None
    }

    if not stored_filename:
        return result

    try:
        # Check if document exists in database using new structure
        db_query = "SELECT original_filename, storage_path FROM documents WHERE stored_filename = %s"
        db_rows, _ = safe_db_query(db_query, (stored_filename,))
        if not isinstance(db_rows, list) or not db_rows:
            db_original_filename = None
            db_storage_path = None
        else:
            db_original_filename = db_rows[0][0] if db_rows[0] else None
            db_storage_path = db_rows[0][1] if len(db_rows[0]) > 1 else None
            result['exists_in_db'] = True  # Document found in database
            result['original_filename'] = db_original_filename or original_filename
            result['storage_path'] = db_storage_path

        # Check if file exists using storage_path from database or fallback to search
        if db_storage_path and os.path.exists(db_storage_path) and os.path.isfile(db_storage_path):
            result['exists_in_file'] = True
            result['file_path'] = db_storage_path
        else:
            # Fallback: Search in the three main directories and user subdirectories
            search_dirs = [
                data_path('documents', 'admin'),
                data_path('documents', 'portal'),
                data_path('documents', 'user'),
            ]
            
            # First check main directories
            for search_dir in search_dirs:
                file_path = os.path.join(search_dir, stored_filename)
                if os.path.exists(file_path) and os.path.isfile(file_path):
                    result['exists_in_file'] = True
                    result['file_path'] = file_path
                    break
            
            # If not found and still looking, check user subdirectories (chat_id folders)
            if not result['exists_in_file']:
                user_dir = data_path('documents', 'user')
                if os.path.exists(user_dir):
                    try:
                        for subdir in os.listdir(user_dir):
                            subdir_path = os.path.join(user_dir, subdir)
                            if os.path.isdir(subdir_path):
                                file_path = os.path.join(subdir_path, stored_filename)
                                if os.path.exists(file_path) and os.path.isfile(file_path):
                                    result['exists_in_file'] = True
                                    result['file_path'] = file_path
                                    break
                    except Exception as e:
                        logging.warning(f"Error searching user subdirectories: {e}")

        # Determine overall status
        if result['exists_in_db'] and result['exists_in_file']:
            result['status'] = 'available'
        elif result['exists_in_db'] and not result['exists_in_file']:
            result['status'] = 'file_missing'
        elif not result['exists_in_db'] and result['exists_in_file']:
            result['status'] = 'db_missing'
        else:
            result['status'] = 'missing'

    except Exception as e:
        logging.warning(f"Error verifying document {stored_filename}: {e}")
        result['status'] = 'error'

    return result

def validate_document_exist_db(stored_filename=None, document_id=None, storage_path=None):
    """
    Simple boolean check if document exists in database.
    Args:
        stored_filename (str, optional): The stored filename (UUID-based) to validate
        document_id (str, optional): Document UUID from documents.id
        storage_path (str, optional): Storage path of the document on disk
    Returns:
        bool: True if document exists in database, False otherwise
    """
    try:
        conditions = []
        params = []

        if stored_filename:
            conditions.append("stored_filename = %s")
            params.append(stored_filename)

        if document_id:
            conditions.append("id = %s")
            params.append(document_id)

        if storage_path:
            conditions.append("storage_path = %s")
            params.append(storage_path)

        if not conditions:
            return False

        where_clause = " OR ".join(conditions)
        query = f"SELECT id FROM documents WHERE {where_clause} LIMIT 1"
        result = safe_db_query(query, tuple(params))

        if isinstance(result, tuple) and len(result) == 2:
            results, columns = result
            if not isinstance(results, list):
                return False
            return bool(results and len(results) > 0)

        return False
    except Exception as e:
        logging.warning(f"Error validating document {stored_filename}: {e}")
        return False

def validate_file_content(content, filename, max_size_mb=50):
    """
    Validate file content before processing.
    Returns (is_valid, reason) tuple.
    """
    try:
        # Check file size
        size_mb = len(content) / (1024 * 1024)
        if size_mb > max_size_mb:
            return False, f"File too large: {size_mb:.1f}MB (max: {max_size_mb}MB)"

        # Check minimum file size (avoid empty or near-empty files)
        if len(content) < 50:
            return False, "File too small (less than 50 bytes)"

        # Check for PDF files
        if filename.lower().endswith('.pdf'):
            # Check PDF magic bytes
            if not content.startswith(b'%PDF-'):
                return False, "Invalid PDF file (missing PDF header)"

            # Check for common HTML error pages disguised as PDFs
            content_str = content.decode('utf-8', errors='ignore').lower()
            if any(html_indicator in content_str for html_indicator in [
                '<!doctype html', '<html>', 'not found', '404 error', 'error page'
            ]):
                return False, "File appears to be an HTML error page, not a valid PDF"

        # Check for text files
        elif filename.lower().endswith(('.txt', '.doc', '.docx')):
            # Try to decode as text to check validity
            try:
                content.decode('utf-8')
            except UnicodeDecodeError:
                try:
                    content.decode('latin-1')
                except UnicodeDecodeError:
                    return False, "Cannot decode text file with UTF-8 or Latin-1 encoding"

        return True, "Valid file"

    except Exception as e:
        return False, f"Validation error: {str(e)}"

def extract_text_from_pdf(path, document_source=None):
    """Extract text from PDF using pdfplumber."""
    if not pdfplumber:
        return ''

    texts = []
    try:
        # Import VisionService only when needed to avoid circular import
        try:
            from app.services.agent.vision_service import VisionService
            vision_service = VisionService()
        except Exception:
            vision_service = None

        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                table_lines = []
                try:
                    tables = page.extract_tables() or []
                except Exception:
                    tables = []
                for table in tables:
                    if not table or len(table) < 2:
                        continue
                    header_index = None
                    header_cells = None
                    for idx, row in enumerate(table):
                        if not row:
                            continue
                        cells = [str(c).strip() if c is not None else "" for c in row]
                        non_empty = [c for c in cells if c]
                        if len(non_empty) >= 2:
                            header_index = idx
                            header_cells = cells
                            break
                    if header_index is None or not header_cells:
                        continue
                    for idx, row in enumerate(table):
                        if idx == header_index or not row:
                            continue
                        values = [str(c).strip() if c is not None else "" for c in row]
                        if not any(values):
                            continue
                        items = []
                        for h, v in zip(header_cells, values):
                            h_clean = str(h).strip() if h is not None else ""
                            if not h_clean or not v:
                                continue
                            items.append({"key": h_clean, "value": v})
                        if items:
                            try:
                                line = json.dumps(items, ensure_ascii=False)
                            except Exception:
                                fallback_parts = [f"{it['key']}: {it['value']}" for it in items]
                                line = " | ".join(fallback_parts)
                            table_lines.append(line)
                if table_lines:
                    texts.extend(table_lines)
                    continue

                text = page.extract_text()
                if text:
                    texts.append(text.strip())
                    continue

                img = page.to_image(resolution=300).original
                img_bytes = None
                try:
                    from io import BytesIO
                    buf = BytesIO()
                    img.save(buf, format="PNG")
                    img_bytes = buf.getvalue()
                except Exception as e_img:
                    logging.warning(f"Failed to convert PDF page to bytes for VisionService: {e_img}")

                text_vision = None
                if vision_service and getattr(vision_service, 'enabled', False) and img_bytes:
                    import base64
                    img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                    try:
                        text_vision = vision_service.describe_image_base64(img_b64, question="Ekstrak semua teks penting dari halaman PDF ini.", mime_type="image/png")
                    except Exception as e_vs:
                        logging.warning(f"VisionService OCR failed: {e_vs}")

                if text_vision and text_vision.strip():
                    texts.append(text_vision.strip())
                else:
                    texts.append(ocr_image(img))

        result = "\n".join(texts)
        if document_source and result:
            logging.info(f"✅ Extracted {len(result)} characters from PDF {document_source}")
        return result
    except Exception as e:
        if document_source:
            logging.warning(f"PDF text extraction failed for {document_source}: {e}")
        return ''

def extract_text_from_docx(file_path, document_source=None):
    """Extract text content from DOCX files."""
    if not docx:
        logging.warning("python-docx not available, skipping DOCX extraction")
        return ''

    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text])
        if document_source and text:
            logging.info(f"✅ Extracted {len(text)} characters from DOCX {document_source}")
        return text
    except Exception as err:
        if document_source:
            logging.warning(f"DOCX extraction failed for {document_source}: {err}")
        return ''

def extract_text_from_xlsx(file_path, document_source=None):
    """Extract text content from XLSX or XLSM files."""
    if not openpyxl:
        logging.warning("openpyxl not available, skipping spreadsheet extraction")
        return ''

    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        texts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            sheet_structured = []
            sheet_simple = []
            headers = None
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell).strip() if cell is not None else "" for cell in row]
                non_empty = [c for c in cells if c]
                if not non_empty:
                    continue
                if headers is None:
                    if len(non_empty) >= 2:
                        headers = cells
                        continue
                    sheet_simple.append(" ".join(non_empty))
                    continue
                items = []
                for h, v in zip(headers, cells):
                    h_clean = str(h).strip() if h is not None else ""
                    if not h_clean or not v:
                        continue
                    items.append({"key": h_clean, "value": v})
                if items:
                    try:
                        line = json.dumps(items, ensure_ascii=False)
                    except Exception:
                        fallback_parts = [f"{it['key']}: {it['value']}" for it in items]
                        line = " | ".join(fallback_parts)
                    sheet_structured.append(line)
            if sheet_structured:
                texts.extend(sheet_structured)
            else:
                texts.extend(sheet_simple)
        text = "\n".join(texts)
        if document_source and text:
            logging.info(f"✅ Extracted {len(text)} characters from spreadsheet {document_source}")
        return text
    except Exception as err:
        if document_source:
            logging.warning(f"Spreadsheet extraction failed for {document_source}: {err}")
        return ''

def extract_text_from_pptx(file_path, document_source=None):
    """Extract text content from PPTX presentations."""
    if not Presentation:
        logging.warning("python-pptx not available, skipping PPTX extraction")
        return ''

    try:
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        text = "\n".join(texts)
        if document_source and text:
            logging.info(f"✅ Extracted {len(text)} characters from presentation {document_source}")
        return text
    except Exception as err:
        if document_source:
            logging.warning(f"PPTX extraction failed for {document_source}: {err}")
        return ''

def extract_text_from_image(path):
    """Extract text from image files using OCR."""
    if not (pytesseract and Image):
        return ''

    try:
        with Image.open(path) as img:
            return ocr_image(img)
    except Exception as e:
        logging.warning(f"Image OCR failed for {path}: {e}")
        return ''

def extract_text_from_image_ocr(file_path, document_source=None, max_pages=3, resolution=150):
    """Backward compatible OCR helper for images and small PDF samples."""
    if not pytesseract:
        logging.warning("pytesseract not available, skipping OCR request")
        return ''

    if not os.path.exists(file_path):
        logging.warning(f"OCR source missing: {file_path}")
        return ''

    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == '.pdf':
        if not pdfplumber:
            return extract_text_from_pdf_image(file_path)

        ocr_texts = []
        try:
            with pdfplumber.open(file_path) as pdf:
                page_limit = min(max_pages, len(pdf.pages))
                for index, page in enumerate(pdf.pages[:page_limit], start=1):
                    try:
                        pil_image = page.to_image(resolution=resolution).original
                        text = ocr_image(pil_image)
                        if text.strip():
                            ocr_texts.append(text.strip())
                    except Exception as page_error:
                        logging.warning(
                            f"OCR failed for page {index} of {document_source or file_path}: {page_error}"
                        )
        except Exception as pdf_error:
            logging.warning(
                f"PDF OCR preprocessing failed for {document_source or file_path}: {pdf_error}"
            )
            return extract_text_from_pdf_image(file_path)

        if ocr_texts:
            return "\n\n".join(ocr_texts)

        return extract_text_from_pdf_image(file_path)

    return extract_text_from_image(file_path)

def extract_document(path):
    """Extract text from various document types - unified function from grabber.py."""
    if not path or not os.path.exists(path):
        return None

    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        try:
            # Try pdfplumber first
            text = extract_text_from_pdf(path, path)
            if text and text.strip():
                return text
        except Exception:
            pass

        try:
            return extract_text_from_pdf_image(path)
        except Exception as e:
            logging.warning(f"PDF extraction failed for {path}: {e}")
            return None

    elif ext == ".docx":
        return extract_text_from_docx(path)
    elif ext in [".xlsx", ".xlsm"]:
        return extract_text_from_xlsx(path)
    elif ext == ".pptx":
        return extract_text_from_pptx(path)
    elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
        return extract_text_from_image(path)
    else:
        return None

def extract_text_from_document(file_path, document_source, use_ocr_fallback=True):
    """
    Extract text from various document types (PDF, text files, images).
    Returns extracted text or empty string if failed.
    Args:
        file_path: Path to the file
        document_source: Source identifier for logging
        use_ocr_fallback: Whether to use OCR if direct extraction fails
    """
    text = ''

    # Check file extension
    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == '.pdf':
        # Try direct PDF text extraction first
        text = extract_text_from_pdf(file_path, document_source)
        logging.info(
            f"✅ Direct PDF extraction success result: {text.strip()}"
        )
        # Use OCR as fallback if no text extracted
        if not text.strip() and use_ocr_fallback:
            file_size = os.path.getsize(file_path)
            size_mb = file_size / (1024 * 1024)
            if file_size < 5 * 1024 * 1024:
                logging.info(
                    f"ℹ️ Direct PDF extraction yielded no text for {document_source}; attempting OCR fallback"
                )
                text = extract_text_from_image_ocr(file_path, document_source)
            else:
                logging.info(
                    f"ℹ️ Skipping OCR fallback for {document_source}; file size {size_mb:.1f}MB exceeds 5MB limit"
                )

            if not text.strip():
                pymupdf_text = extract_text_from_pdf_image(
                        file_path
                    )
                if pymupdf_text.strip():
                    logging.info(
                        f"ℹ️ PyMuPDF OCR fallback extracted {len(pymupdf_text)} characters for {document_source}"
                    )
                    text = pymupdf_text

    elif file_ext == '.docx':
        text = extract_text_from_docx(file_path, document_source)

    elif file_ext in ['.xlsx', '.xlsm']:
        text = extract_text_from_xlsx(file_path, document_source)

    elif file_ext == '.pptx':
        text = extract_text_from_pptx(file_path, document_source)

    elif file_ext in ['.txt', '.md', '.log']:
        # For text files, try to read directly
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            logging.info(f"✅ Read {len(text)} characters from text file {document_source}")
        except Exception as e:
            try:
                # Try with different encoding
                with open(file_path, 'r', encoding='latin-1') as f:
                    text = f.read()
                logging.info(f"✅ Read {len(text)} characters from text file {document_source} (latin-1 encoding)")
            except Exception as e2:
                logging.warning(f"Failed to read text file {document_source}: {e} (utf-8), {e2} (latin-1)")

    elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif'] and use_ocr_fallback:
        # For image files, use OCR directly
        text = extract_text_from_image(file_path)

    else:
        logging.warning(f"Unsupported file type: {file_ext} for {document_source}")

    return text

def process_document_for_vector_storage(
    file_path,
    document_name,
    document_source,
    metadata=None,
    document_id: Optional[str] = None,
    storage_path: Optional[str] = None
):
    """
    Process a document (PDF, text, or image) and add it to vector storage.
    Returns True if successful, False otherwise.
    """
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        from langchain_core.documents import Document

        vectorstore = get_vectorstore()
        if not vectorstore:
            logging.warning("❌ Vectorstore not available for document processing")
            return False

        file_path_resolved = os.path.abspath(file_path)
        if not os.path.exists(file_path_resolved):
            logging.warning(f"⚠️ File not found for vector processing: {file_path}")
            return False

        text = extract_text_from_document(file_path_resolved, document_source)
        if not text.strip():
            logging.warning(f"⚠️ No text extracted from {document_source}, skipping vector storage")
            return False

        db_record = None
        if not document_id or not storage_path:
            try:
                query = """
                    SELECT id, storage_path, mime_type, source_type, uploaded_by
                    FROM documents
                    WHERE stored_filename = %s
                    LIMIT 1
                """
                result, columns = safe_db_query(query, (document_source,))
                if isinstance(result, list) and result:
                    db_record = dict(zip(columns, result[0]))
            except Exception as db_err:
                logging.warning(f"Error fetching document record for {document_source}: {db_err}")

        if db_record:
            document_id = document_id or db_record.get('id')
            storage_path = storage_path or db_record.get('storage_path')
            mime_type = (metadata or {}).get('mime_type') or db_record.get('mime_type')
            source_type = (metadata or {}).get('source_type') or db_record.get('source_type')
            uploaded_by = (metadata or {}).get('uploaded_by') or db_record.get('uploaded_by')
        else:
            mime_type = (metadata or {}).get('mime_type')
            source_type = (metadata or {}).get('source_type')
            uploaded_by = (metadata or {}).get('uploaded_by')

        if not document_id:
            logging.warning(f"⚠️ Document ID missing for {document_source}, aborting vector storage")
            return False

        if not storage_path:
            storage_path = os.path.relpath(file_path_resolved, '.')

        if not mime_type:
            mime_type, _ = mimetypes.guess_type(document_name)
            if not mime_type:
                mime_type = 'application/octet-stream'

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        raw_chunks = text_splitter.split_text(text)
        chunks = [chunk for chunk in raw_chunks if chunk.strip()]

        if not chunks:
            logging.warning(f"⚠️ No valid chunks created for {document_source}")
            return False

        chunk_total = len(chunks)
        created_at = get_current_datetime().isoformat()

        base_metadata = {
            "document_id": str(document_id),
            "document_source": document_source,
            "stored_filename": document_source,
            "document_name": document_name,
            "original_filename": document_name,
            "source_type": source_type,
            "storage_path": storage_path,
            "source": storage_path,
            "mime_type": mime_type,
            "uploaded_by": uploaded_by,
            "chunk_total": chunk_total,
            "created_at": created_at,
        }

        if metadata:
            base_metadata.update(metadata)

        docs = []
        display_name = document_name or os.path.basename(str(document_source or "")) or ""
        prefix = f"{display_name}\n\n" if display_name else ""
        for index, chunk in enumerate(chunks):
            doc_metadata = base_metadata.copy()
            doc_metadata["chunk_index"] = index
            content = f"{prefix}{chunk}" if prefix else chunk
            docs.append(Document(page_content=content, metadata=doc_metadata))

        try:
            vectorstore.add_documents(docs)
            logging.info(f"✅ Added {len(docs)} chunks to vector store for {document_source}")
            return True
        except Exception as add_err:
            logging.error(f"❌ Failed to add chunks to vector store for {document_source}: {add_err}")
        return False

    except Exception as e:
        logging.error(f"❌ Failed to process document {document_source} for vector storage: {e}")
        return False


def cleanup_orphan_document_files() -> dict:
    """
    Remove files under data/documents for source_type 'portal' or 'website'
    that no longer exist in the documents table. Only deletes related files
    in their respective directories and logs actions.

    Returns a summary dict with counts and details.
    """
    summary = {
        "checked": 0,
        "deleted": 0,
        "kept": 0,
        "errors": []
    }

    try:
        rows, cols = safe_db_query(
            "SELECT stored_filename, storage_path, source_type FROM documents WHERE source_type IN ('portal','website')"
        )
        valid_files = set()
        for r in rows or []:
            stored = r[0]
            storage_path = r[1] if len(r) > 1 else None
            source_type = r[2] if len(r) > 2 else None
            if stored:
                valid_files.add((source_type or '', stored))
                # Also trust absolute storage_path if present
                if storage_path and os.path.isfile(storage_path):
                    valid_files.add((source_type or '', os.path.basename(storage_path)))

        targets = [
            ("portal", data_path("documents", "portal")),
            ("website", data_path("documents", "website")),
        ]

        for src_type, folder in targets:
            if not os.path.isdir(folder):
                continue
            for name in os.listdir(folder):
                path = os.path.join(folder, name)
                # only consider regular files; skip subdirs and hidden temp files
                if not os.path.isfile(path) or name.startswith('.'):
                    continue
                summary["checked"] += 1
                if (src_type, name) in valid_files:
                    summary["kept"] += 1
                    continue
                try:
                    os.remove(path)
                    summary["deleted"] += 1
                    logging.info(f"🧹 Deleted orphan document file: {path}")
                except Exception as err:
                    msg = f"Failed to delete {path}: {err}"
                    logging.warning(msg)
                    summary["errors"].append(msg)

    except Exception as e:
        msg = f"Cleanup failed: {e}"
        logging.error(msg)
        summary["errors"].append(msg)

    logging.info(
        f"🧾 Cleanup summary — checked: {summary['checked']}, kept: {summary['kept']}, deleted: {summary['deleted']}, errors: {len(summary['errors'])}"
    )
    return summary
